#!/usr/bin/env python3
import os
import sys
import re
from collections import defaultdict

from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.enum.shapes import MSO_CONNECTOR

# --- Helper function to parse a filename ---
def parse_filename(filename):
    """
    Expected filename format (underscore-separated parts):
      Part0: arbitrary (e.g. "C4-MAX")
      Part1: brain region and section, e.g. "DG1" (region = DG, section = 1)
      Part2: image number (e.g. "1")
      Part3: animal+genotype+treatment+sex, e.g. "G190E4FADLPSFEMALE"
      The rest of the parts are ignored.
    Returns a dictionary with keys:
      'brain_region', 'section', 'image_number', 'animal', 'genotype', 'treatment', 'sex'
    """
    name, ext = os.path.splitext(filename)
    parts = name.split('_')
    if len(parts) < 4:
        return None  # skip if not enough parts

    # Part 1: extract brain region (letters) and section (digits)
    m = re.match(r"([A-Za-z]+)(\d+)", parts[1])
    if not m:
        return None
    brain_region = m.group(1)
    section = m.group(2)

    image_number = parts[2]

    # Part 3: animal identifier + genotype + treatment + sex.
    if len(parts) >= 5:
        animal = parts[3]
        info = parts[4]
        rest = info
    else:
        info = parts[3]
        m_animal = re.match(r"^([A-Z]\d+)", info)
        if not m_animal:
            return None
        animal = m_animal.group(1)
        rest = info[len(animal):]  # remove animal id

    # Determine treatment by checking for 'LPS' or 'VEH'
    if "LPS" in rest:
        treatment = "LPS"
    elif "VEH" in rest:
        treatment = "VEH"
    else:
        treatment = "UNKNOWN"

    # Split rest by treatment string to get genotype and sex.
    split_parts = re.split(r"(LPS|VEH)", rest)
    if len(split_parts) < 3:
        return None
    genotype = split_parts[0]
    sex = split_parts[2]

    return {
        'brain_region': brain_region,
        'section': section,
        'image_number': image_number,
        'animal': animal,
        'genotype': genotype,
        'treatment': treatment,
        'sex': sex,
        'filename': filename
    }

# --- Main function ---
def main():
    # Process command-line arguments:
    if len(sys.argv) > 1:
        image_dir = sys.argv[1]
    else:
        print("Usage: {} <image_directory> [output_file]".format(sys.argv[0]))
        sys.exit(1)
    if len(sys.argv) > 2:
        output_file = sys.argv[2]
    else:
        output_file = 'output.pptx'

    # List image files in the directory (filter jpg, jpeg, png)
    valid_exts = ('.jpg', '.jpeg', '.png')
    all_files = os.listdir(image_dir)
    image_files = [f for f in all_files if f.lower().endswith(valid_exts)]
    if not image_files:
        print("No valid image files found in directory:", image_dir)
        sys.exit(1)

    # Parse filenames and collect info.
    images_data = []
    for f in image_files:
        parsed = parse_filename(f)
        if parsed:
            parsed['filepath'] = os.path.join(image_dir, f)
            images_data.append(parsed)
        else:
            print("Skipping file (could not parse):", f)

    # Group images by animal, genotype, treatment
    groups = defaultdict(list)
    for data in images_data:
        key = (data['animal'], data['genotype'], data['treatment'])
        groups[key].append(data)

    # Create presentation
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # blank slide

    # Layout constants (using inches for arithmetic)
    title_top = 0.5
    title_height = 0.5
    row_spacing = 0.3
    img_label_height = 0.2
    bracket_line_height = 0.05
    section_label_height = 0.2
    cm_to_in = 0.3937

    # Process each group to create one slide per group.
    for key, group_images in groups.items():
        animal, genotype, treatment = key
        total_images = len(group_images)
        size_cm = 7.21 if total_images < 9 else 6.21
        image_size = size_cm * cm_to_in  # in inches

        # Create a new slide with a title that does not include brain region.
        slide = prs.slides.add_slide(blank_slide_layout)
        title_text = "Genotype: {} | Treatment: {} | Animal: {}".format(genotype, treatment, animal)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(title_top), prs.slide_width - Inches(1), Inches(title_height))
        title_box.text_frame.text = title_text

        # Group images by section (with brain region) for this slide.
        section_groups = defaultdict(list)
        for data in group_images:
            section_key = (data['brain_region'], data['section'])
            section_groups[section_key].append(data)

        # Sort sections by section number (ascending)
        sorted_sections = sorted(section_groups.items(),
                                 key=lambda x: int(x[0][1]) if x[0][1].isdigit() else x[0][1])
        
        # Starting vertical position (below title)
        current_top = title_top + title_height + 0.3

        # Layout each section row.
        for (brain_region, section), imgs in sorted_sections:
            imgs_sorted = sorted(imgs, key=lambda d: int(d['image_number']) if d['image_number'].isdigit() else d['image_number'])
            n_imgs = len(imgs_sorted)
            spacing_x = 0.2  # inches (raw)
            total_row_width = n_imgs * image_size + (n_imgs - 1) * spacing_x
            slide_width_in = prs.slide_width.inches
            start_left = (slide_width_in - total_row_width) / 2

            # Add images and their labels.
            for i, data in enumerate(imgs_sorted):
                left = Inches(start_left + i * (image_size + spacing_x))
                pic = slide.shapes.add_picture(data['filepath'], left, Inches(current_top), width=Inches(image_size), height=Inches(image_size))
                label_top = current_top + image_size
                label_box = slide.shapes.add_textbox(left, Inches(label_top), Inches(image_size), Inches(img_label_height))
                label_box.text_frame.text = "Img " + data['image_number']
                label_box.text_frame.paragraphs[0].font.size = Pt(10)

            # Add a horizontal line (bracket) below the image labels.
            line_top = current_top + image_size + img_label_height + 0.05
            line_left = Inches(start_left)
            # Note: Convert the end point properly: start_left + total_row_width (in inches) then to Inches()
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                line_left, Inches(line_top),
                Inches(start_left + total_row_width), Inches(line_top)
            )

            # Add section label (with brain region) below the line.
            section_label_top = line_top + 0.05
            sec_label_box = slide.shapes.add_textbox(Inches(start_left), Inches(section_label_top), Inches(total_row_width), Inches(section_label_height))
            sec_label_box.text_frame.text = "{} Section {}".format(brain_region, section)
            for para in sec_label_box.text_frame.paragraphs:
                para.font.size = Pt(10)
                para.alignment = 1  # centered

            # Update vertical position for next section row.
            current_top = section_label_top + section_label_height + row_spacing

    prs.save(output_file)
    print("Presentation saved as", output_file)

if __name__ == '__main__':
    main()